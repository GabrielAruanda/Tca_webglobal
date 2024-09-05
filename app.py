from flask import Flask, render_template, request, redirect, url_for, send_file
# Flask: Um framework web para Python que facilita a criação de aplicativos web.
# render_template: Função para renderizar templates HTML.
# request: Objeto que contém os dados da solicitação HTTP.
# redirect: Função para redirecionar o usuário para uma URL diferente.
# url_for: Função para gerar URLs para as funções de view.
# send_file: Função para enviar arquivos como resposta HTTP.

import mysql.connector
# mysql.connector: Conector para se conectar e interagir com um banco de dados MySQL.

import pandas as pd
# pandas: Biblioteca para manipulação e análise de dados, especialmente útil para operações em DataFrames e leitura/escrita de arquivos.

from docx import Document
# docx: Biblioteca para criar e manipular documentos no formato DOCX (Microsoft Word).

from pptx import Presentation
from pptx.util import Inches
# pptx: Biblioteca para criar e manipular apresentações no formato PPTX (Microsoft PowerPoint).
# Inches: Utilitário para definir tamanhos e posições em polegadas para os elementos dos slides.

import io
# io: Biblioteca para manipulação de fluxo de dados em memória, como leitura e escrita de arquivos em bytes.

import hashlib
# hashlib: Biblioteca para criar hashes criptográficos, usada para gerar códigos curtos únicos.

from datetime import datetime
# datetime: Biblioteca para manipulação de datas e horários.

import plotly.express as px
import plotly.io as pio
# plotly.express: Biblioteca para criação de gráficos interativos de forma simplificada.
# plotly.io: Biblioteca para entrada/saída de gráficos, incluindo exportação de imagens.

app = Flask(__name__)
# Cria uma instância da aplicação Flask.

# Configuração do banco de dados
db_config = {
    'user': 'root',                # Nome de usuário do banco de dados MySQL.
    'password': '',                # Senha do banco de dados MySQL (vazia neste caso).
    'host': 'localhost',           # Endereço do servidor MySQL.
    'database': 'tca'              # Nome do banco de dados.
}

db = mysql.connector.connect(**db_config)
# Estabelece a conexão com o banco de dados MySQL usando as configurações fornecidas.

def generate_unique_short_code():
    """Gera um código curto único."""
    while True:
        # Gera um código curto usando o hash MD5 da data e hora atual.
        short_code = hashlib.md5(datetime.now().strftime('%Y-%m-%d %H:%M:%S').encode()).hexdigest()[:6]
        cursor = db.cursor()
        # Verifica no banco de dados se o código já existe.
        cursor.execute("SELECT COUNT(*) FROM urls WHERE short_code = %s", (short_code,))
        if cursor.fetchone()[0] == 0:
            cursor.close()
            return short_code
        cursor.close()

@app.route("/", methods=["GET", "POST"])
def index():
    """Página principal para encurtar URLs."""
    if request.method == "POST":
        # Recebe a URL original do formulário.
        original_url = request.form["original_url"]
        # Gera um código curto único para a URL.
        short_code = generate_unique_short_code()

        cursor = db.cursor()
        # Insere a URL original, o código curto, a data de criação e a contagem de cliques inicial no banco de dados.
        cursor.execute("INSERT INTO urls (original_url, short_code, created_at, click_count) VALUES (%s, %s, %s, %s)",
                       (original_url, short_code, datetime.now(), 0))
        db.commit()
        cursor.close()

        # Gera a URL curta completa para redirecionar.
        short_url = url_for("redirect_url", short_code=short_code, _external=True)
        return render_template("index.html", short_url=short_url)

    return render_template("index.html")

@app.route("/<short_code>")
def redirect_url(short_code):
    """Redireciona para a URL original com base no código curto."""
    cursor = db.cursor(dictionary=True)
    # Recupera a URL original usando o código curto.
    cursor.execute("SELECT original_url FROM urls WHERE short_code = %s", (short_code,))
    url_data = cursor.fetchone()

    if url_data:
        client_ip = request.remote_addr
        # Atualiza o registro com a data e IP do último clique e incrementa a contagem de cliques.
        cursor.execute("""
            UPDATE urls 
            SET last_click_at = %s, last_click_ip = %s, click_count = click_count + 1 
            WHERE short_code = %s
        """, (datetime.now(), client_ip, short_code))
        db.commit()
        cursor.close()
        # Redireciona o usuário para a URL original.
        return redirect(url_data['original_url'])
    else:
        cursor.close()
        return "URL não encontrada", 404

@app.route("/urls")
def show_urls():
    """Exibe todas as URLs encurtadas e suas informações."""
    cursor = db.cursor(dictionary=True)
    # Recupera todas as URLs do banco de dados.
    cursor.execute("SELECT * FROM urls")
    urls = cursor.fetchall()
    cursor.close()
    return render_template("urls.html", urls=urls)

@app.route("/charts")
def charts():
    """Exibe um gráfico com a contagem de cliques para cada URL."""
    cursor = db.cursor(dictionary=True)
    # Recupera dados das URLs do banco de dados.
    cursor.execute("SELECT short_code, original_url, click_count, created_at, last_click_at, last_click_ip FROM urls")
    urls = cursor.fetchall()
    cursor.close()

    # Preparar os dados para o gráfico.
    df = pd.DataFrame(urls)
    fig = px.bar(df, x='short_code', y='click_count', labels={'short_code': 'Short Code', 'click_count': 'Clicks'}, title='URL Clicks')

    # Converter gráfico para HTML.
    chart_html = fig.to_html(full_html=False)

    return render_template("charts.html", url_data=urls, chart_html=chart_html)

@app.route("/download/<file_type>")
def download_report(file_type):
    """Gera e baixa um relatório no formato especificado (pptx, docx, xlsx)."""
    cursor = db.cursor(dictionary=True)
    # Recupera dados das URLs do banco de dados.
    cursor.execute("SELECT short_code, original_url, click_count, created_at, last_click_at, last_click_ip FROM urls")
    urls = cursor.fetchall()
    cursor.close()

    if file_type == 'pptx':
        # Cria uma apresentação PowerPoint.
        prs = Presentation()
        
        # Gerar o gráfico e adicionar ao PowerPoint.
        df = pd.DataFrame(urls)
        fig = px.bar(df, x='short_code', y='click_count', labels={'short_code': 'Short Code', 'click_count': 'Clicks'}, title='URL Clicks')
        img_bytes = pio.to_image(fig, format='png')
        
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "URL Clicks"

        # Adicionar gráfico ao slide.
        image_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(image_stream, Inches(1), Inches(1.5), width=Inches(8))

        # Adicionar informações das URLs.
        for url in urls:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = f"Short Code: {url['short_code']}"

            content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
            text_frame = content.text_frame
            p = text_frame.add_paragraph()
            p.text = (
                f"Original URL: {url['original_url']}\n"
                f"Clicks: {url['click_count']}\n"
                f"Created At: {url['created_at'].strftime('%Y-%m-%d %H:%M:%S')}\n"
                f"Last Click At: {url['last_click_at'].strftime('%Y-%m-%d %H:%M:%S') if url['last_click_at'] else 'Never'}\n"
                f"Last Click IP: {url['last_click_ip'] if url['last_click_ip'] else 'N/A'}"
            )

        # Salvar o arquivo pptx em memória.
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        return send_file(file_stream, as_attachment=True, download_name="report.pptx", mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    elif file_type == 'docx':
        # Cria um documento Word.
        doc = Document()
        doc.add_heading('URL Report', 0)

        for url in urls:
            doc.add_paragraph(f"Short Code: {url['short_code']}")
            doc.add_paragraph(f"Original URL: {url['original_url']}")
            doc.add_paragraph(f"Clicks: {url['click_count']}")
            doc.add_paragraph(f"Created At: {url['created_at'].strftime('%Y-%m-%d %H:%M:%S')}")
            doc.add_paragraph(f"Last Click At: {url['last_click_at'].strftime('%Y-%m-%d %H:%M:%S') if url['last_click_at'] else 'Never'}")
            doc.add_paragraph(f"Last Click IP: {url['last_click_ip'] if url['last_click_ip'] else 'N/A'}")
            doc.add_paragraph('')  # Linha em branco para separar as URLs.

        # Salvar o arquivo docx em memória.
        file_stream = io.BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        return send_file(file_stream, as_attachment=True, download_name="report.docx", mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

    elif file_type == 'xlsx':
        # Cria um DataFrame e converte para um arquivo Excel.
        df = pd.DataFrame(urls)
        df['created_at'] = df['created_at'].apply(lambda x: x.strftime('%Y-%m-%d %H:%M:%S'))
        df['last_click_at'] = df['last_click_at'].apply(lambda x: x.strftime('%Y-%m-%d %H:%M:%S') if x else 'Never')
        df['last_click_ip'] = df['last_click_ip'].fillna('N/A')

        file_stream = io.BytesIO()
        df.to_excel(file_stream, index=False)
        file_stream.seek(0)

        return send_file(file_stream, as_attachment=True, download_name="report.xlsx", mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

    else:
        return "Tipo de arquivo não suportado", 400

if __name__ == "__main__":
    app.run(debug=True)
# Executa a aplicação Flask em modo de depuração.
